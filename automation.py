# Script: full_automation_pipeline.py
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import os
import os
import re
from pptx import Presentation
from pptx.util import Pt

# === CONFIGURATION ===

# Paths to theme pivot tables (quantitative)
pivot_paths = {
    "Youth": "src/data/Youth - jeunesse.xlsx",
    "Education": "src/data/Education - éducation.xlsx",
    "Digital": "src/data/Digital - numérique.xlsx",
    "IFF & Transboundary Crime": "src/data/IFF - Transnational _Crimes.xlsx",
    "Mining": "src/data/Mining - Mine.xlsx"
}

pivot_sheet_names = {
    "Youth": "Youth - Pivot Table",
    "Education": "Education - Pivote Table",
    "Digital": "Digital - Pivot Table",
    "IFF & Transboundary Crime": "IFF & Crime - Pivot Table",
    "Mining": "Mining - Pivot Table"
}

# Paths to qualitative summaries
summary_files = {
    "Youth": "src/o1_outputs/Youth_Output.txt",
    "Education": "src/o1_outputs/Education_Output.txt",
    "Digital": "src/o1_outputs/Digital_Output.txt",
    "IFF & Transboundary Crime": "src/o1_outputs/IFF - Transnational_Crimes_Output.txt",
    "Mining": "src/o1_outputs/Mining_Output.txt"
}

# Output files
pptx_input_path = "JWP.pptx"
pptx_output_path = "src/data/Thematic_Areas_Updated.pptx"
thinkcell_excel_path = "src/data/ThinkCell_Data.xlsx"

# Load summaries into memory
summaries = {}
for theme, path in summary_files.items():
    with open(path, 'r', encoding='utf-8') as f:
        summaries[theme] = f.read()

# Helper to split summary into 4 focus areas
def segment_summary(text):
    segments = re.split(r"\n\s*\d[\)\.]|(?=\n\s*- )", text.strip())
    segments = [s.strip("-• \n") for s in segments if s.strip()]
    return segments[:4]

# Extract funding and support details from summary
def extract_funding_info(text):
    funding_gap = re.search(r"funding gap.*?(\d{1,3}) ?%", text, re.IGNORECASE)
    required_funding = re.search(r"required funding.*?\$?US?\$?[\s]*([0-9,.]+)\s*(million|billion)", text, re.IGNORECASE)
    funding_situation = re.search(r"(Funding situation.*?)\n", text, re.IGNORECASE | re.DOTALL)
    main_areas = re.search(r"(Main areas of support.*?)\n", text, re.IGNORECASE | re.DOTALL)

    return {
        "funding_gap": funding_gap.group(1) + "%" if funding_gap else None,
        "required_funding": required_funding.group(0) if required_funding else None,
        "funding_situation": funding_situation.group(1) if funding_situation else None,
        "main_areas": main_areas.group(1) if main_areas else None
    }

# Match slide title to theme key
def match_theme(title):
    title = title.lower()
    for theme in summaries:
        if theme in title:
            return theme
    return None

# Load PowerPoint
prs = Presentation(pptx_input_path)

# Loop through slides and fill them
for slide in prs.slides:
    title_shape = next((s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()), None)
    if not title_shape:
        continue

    matched_theme = match_theme(title_shape.text)
    if not matched_theme:
        continue

    segments = segment_summary(summaries[matched_theme])
    funding_info = extract_funding_info(summaries[matched_theme])

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        placeholder = shape.text_frame.text.strip().lower()
        shape.text_frame.clear()
        p = shape.text_frame.paragraphs[0]
        p.font.size = Pt(12)

        # Fill in focus sections
        if "focus 1" in placeholder and len(segments) > 0:
            p.text = segments[0]
        elif "focus 2" in placeholder and len(segments) > 1:
            p.text = segments[1]
        elif "focus 3" in placeholder and len(segments) > 2:
            p.text = segments[2]
        elif "focus 4" in placeholder and len(segments) > 3:
            p.text = segments[3]
        elif "funding gap" in placeholder and funding_info["funding_gap"]:
            p.text = f"Funding gap: {funding_info['funding_gap']}"
        elif "required funding" in placeholder and funding_info["required_funding"]:
            p.text = funding_info["required_funding"]
        elif "funding situation" in placeholder and funding_info["funding_situation"]:
            p.text = funding_info["funding_situation"]
        elif "main areas of support" in placeholder and funding_info["main_areas"]:
            p.text = funding_info["main_areas"]

# Save final result
prs.save(pptx_output_path)
print(f"Saved final presentation to: {pptx_output_path}")